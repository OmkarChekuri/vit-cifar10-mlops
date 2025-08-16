from pptx import Presentation
from pptx.util import Inches, Pt

# Create a new presentation
prs = Presentation()

# Define slide layout (0 = title, 1 = title+content, etc.)
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Vision Transformer for CIFAR-10: An MLOps Approach"
subtitle.text = "A Project Report on a Reproducible Deep Learning Pipeline\n\nYour Name / Project ID: [Your Name]\nDate: [Current Date]"

# Slide 2: Executive Summary
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Executive Summary"
content.text = (
    "Objective: To fine-tune a compact Vision Transformer (vit_tiny_patch16_224) for image classification on the CIFAR-10 dataset.\n\n"
    "Achievement: Achieved a 96.67% test accuracy after 3 epochs of training.\n\n"
    "Impact: Developed a complete, reproducible MLOps pipeline using Docker and GitHub Actions, ensuring consistent performance and a solid foundation for future projects.\n\n"
    "Key Takeaway: Transfer learning with efficient architectures is highly effective and scalable, even on consumer-grade hardware."
)

# Slide 3: Problem & Methodology
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Problem & Methodology"
content.text = (
    "Problem: Building a reliable and efficient image classifier for a standard benchmark.\n\n"
    "Dataset: CIFAR-10 (60,000 images, 10 classes). Data was split into a training set (40k), a validation set (10k), and a test set (10k).\n\n"
    "Approach:\n"
    "- Model: Pre-trained vit_tiny_patch16_224 from the timm library.\n"
    "- Training: AdamW optimizer, Cosine Annealing LR scheduler, and mixed-precision (bfloat16) for efficiency.\n"
    "- Evaluation: Monitored performance on a validation set to save the best model. Final evaluation on a held-out test set."
)

# Slide 4: Performance Results: Loss Curve
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Performance Results: Loss Curve"
content.text = (
    "Training and Validation Loss over 3 Epochs\n\n"
    "Key Points:\n"
    "- The training loss consistently decreased with each epoch, indicating the model was learning effectively.\n"
    "- The validation loss also steadily decreased, showing that the model was generalizing well to unseen data.\n"
    "- This downward trend for both metrics suggests the model is not overfitting.\n\n"
    "[Insert the plot of Training and Validation Loss]"
)

# Slide 5: Performance Results: Accuracy Curve
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Performance Results: Accuracy Curve"
content.text = (
    "Training and Validation Accuracy over 3 Epochs\n\n"
    "Key Points:\n"
    "- Training accuracy steadily increased, reaching almost 90% by the end.\n"
    "- Validation accuracy showed a strong, consistent increase, ending at 89.33%.\n"
    "- The consistent performance gain confirms the efficacy of the transfer learning approach on this dataset.\n\n"
    "[Insert the plot of Training and Validation Accuracy]"
)

# Insert plot into Slide 5 (Accuracy Curve)
left = Inches(5.5)
top = Inches(2)
height = Inches(3)

slide.shapes.add_picture("./Reports/cifar10/accuracy_curve.png", left, top, height=height)


# Slide 6: MLOps & Reproducibility
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "MLOps & Reproducibility"
content.text = (
    "Challenge: Ensured the project was fully reproducible, addressing common issues that arise in development.\n\n"
    "Solutions:\n"
    "- Docker: Containerized the entire environment with a Dockerfile, guaranteeing that the code runs the same way on any machine. Solved shared memory errors with the --shm-size flag.\n"
    "- CI/CD: Implemented a GitHub Actions workflow to automatically test the training script and lint the code on every push, maintaining code quality.\n"
    "- Logging: Integrated a robust logging system to capture all metrics, creating an auditable record for every training run."
)

# Slide 7: Key Challenges & Lessons Learned
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Key Challenges & Lessons Learned"
content.text = (
    "1. Version Mismatches: A version conflict between the Docker base image's CUDA version and PyTorch's required CUDA version was a major roadblock.\n"
    "   Lesson: Explicitly matching library versions with the environment is crucial for reproducibility.\n\n"
    "2. Shared Memory: The DataLoader failed in Docker due to insufficient shared memory.\n"
    "   Lesson: The docker run --shm-size flag is necessary for multi-worker tasks.\n\n"
    "3. Library Bugs: An OverflowError was traced to a bug in a specific version of torchvision's ColorJitter transformation.\n"
    "   Lesson: Constant testing and being aware of library-specific bugs are essential."
)

# Slide 8: Future Work
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Future Work"
content.text = (
    "- Optimization: Explore different hyperparameters and advanced augmentation techniques to push performance even further.\n\n"
    "- Multi-View Project: Complete the training and analysis for the multi-view adapter project and compare its performance to the standard fine-tuning approach.\n\n"
    "- Production Deployment: Build a full deployment pipeline using the provided Jenkinsfile to serve the trained model as a web service."
)

# Save presentation
file_path = "./Reports/cifar10/Vision_Transformer_CIFAR10_Presentation.pptx"
prs.save(file_path)

file_path
