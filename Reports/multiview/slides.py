from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

# Layouts
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]

# Slide 1: Title Slide
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Multi-View Classification: A Memory-Efficient Approach with Adapters"
subtitle.text = "A Project Report on Parameter-Efficient Fine-Tuning (PEFT)\n\nYour Name / Project ID: [Your Name]\nDate: [Current Date]"

# Slide 2: Executive Summary
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Executive Summary"
content.text = (
    "Objective: To develop a memory-efficient solution for multi-view image classification using Parameter-Efficient Fine-Tuning (PEFT) with a Vision Transformer.\n\n"
    "Achievement: Successfully trained a model with only 5.95% of its total parameters being trainable. Achieved a mean validation accuracy of 61.7% across three simulated folds.\n\n"
    "Impact: Validated a highly efficient method for fine-tuning large models on limited hardware, demonstrating a viable approach for medical imaging and other complex tasks."
)

# Slide 3: Problem & Methodology
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Problem & Methodology"
content.text = (
    "Problem: Training multi-view classifiers on large pre-trained models can be memory-intensive and computationally expensive.\n\n"
    "Approach:\n"
    "- Model: Pre-trained vit_tiny_patch16_224 backbone.\n"
    "- PEFT Strategy: Froze the majority of the model and injected small, trainable Adapter modules into its layers.\n"
    "- Data: Used a synthetic DummyMultiViewDataset with four image views for a binary classification task.\n\n"
    "Training Details:\n"
    "- A 3-fold cross-validation was simulated to assess model stability.\n"
    "- Only 0.362 million of the 6.1 million parameters were fine-tuned."
)

# Slide 4: Performance Results
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Performance Results: Loss & Accuracy"
content.text = (
    "Training and Validation Performance Across Three Folds\n\n"
    "Key Points:\n"
    "- The learning curves show a stable, if modest, learning trend across all folds.\n"
    "- The model's performance on a random dataset is limited but consistent, validating the architectural approach.\n\n"
    "[Insert the plots of Training and Validation Loss and Accuracy per Fold]"
)

left = Inches(5.5)
top = Inches(2)
height = Inches(3)

slide.shapes.add_picture("./Reports/multiview/multi_view_performance.png", left, top + Inches(3.5), height=height)

# Slide 5: Key Challenges & Lessons Learned
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Key Challenges & Lessons Learned"
content.text = (
    "1. Memory Constraints: The project successfully demonstrated that adapter-based PEFT is a powerful way to reduce the number of trainable parameters and manage memory.\n\n"
    "2. Dummy Data: The use of synthetic data, while great for architectural testing, showed that a real-world dataset is required to achieve meaningful accuracy gains.\n\n"
    "3. Training Stability: A small batch size (1) was necessary due to the multi-view input, which can lead to noisy gradients and less stable training."
)

# Slide 6: MLOps & Reproducibility
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "MLOps & Reproducibility"
content.text = (
    "Containerization: The Dockerfile provides a consistent environment. We used the docker run --shm-size flag to handle multi-worker data loading, which is a key learning for parallel processing in containers.\n\n"
    "Version Control: All architectural changes and MLOps scripts are versioned with Git.\n\n"
    "Logging: Detailed per-epoch metrics for each fold were logged, providing a comprehensive record for future analysis."
)

# Slide 7: Future Work & Recommendations
slide = prs.slides.add_slide(content_slide_layout)
title, content = slide.shapes.title, slide.placeholders[1]
title.text = "Future Work & Recommendations"
content.text = (
    "- Real-World Data: The next logical step is to train this model on a real multi-view dataset, such as medical scans, to evaluate its true performance.\n\n"
    "- Hyperparameter Tuning: Systematically tune the adapter bottleneck_ratio and learning rate to optimize the model for real data.\n\n"
    "- Model Comparison: Compare the adapter-based approach to a full fine-tuning pipeline to quantify the trade-off between performance and efficiency."
)

# Save file
prs.save("./Reports/multiview/Multi_View_Adapters_Presentation.pptx")
print("Presentation saved as Multi_View_Adapters_Presentation.pptx")
